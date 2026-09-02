import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def apply_text_styling(paragraph, text, font_size, color, bold=False, italic=False, font_name="Arial", align=PP_ALIGN.LEFT):
    paragraph.text = text
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.italic = italic
    paragraph.font.color.rgb = color
    paragraph.alignment = align

def add_clean_container(slide, left, top, width, height, bg_color, border_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

def create_editorial_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # High-End Human Aesthetic Editorial Palette
    COLOR_BG = RGBColor(0xF7, 0xF5, 0xF0)         # Warm Alabaster / Beige Base
    COLOR_NAVY = RGBColor(0x1E, 0x29, 0x3B)       # Deep Slate Navy (Structural Contrast)
    COLOR_COPPER = RGBColor(0xB8, 0x5A, 0x38)     # Burnt Terracotta / Copper Accent
    COLOR_MUTED = RGBColor(0x64, 0x74, 0x8B)      # Editorial Ash Gray Text
    COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)      # Text contrast inside navy blocks

    slides_content = [
        {
            "num": "01",
            "title": "STARTNERVE INTELLIGENCE",
            "subtitle": "THE REGULATORY FIREWALL FOR BULK API MANUFACTURING",
            "left_head": "TITAN V11 SCREENING FRAMEWORK",
            "left_body": "A hybrid dual-stream graph neural network architected for predictive ICH M7 mutagenicity screening. Processing structural coordinates under a minute per compound.",
            "right_head": "VERIFIED DATA MOAT METRICS",
            "right_body": "58,013 fully geometry-resolved 3D molecular graphs. Achieved a rock-solid 0.8603 ROC-AUC on the p53 pathway under strict Bemis-Murcko scaffold splitting.",
            "notes": "Every week a bulk API batch sits waiting on an Ames assay is a week of regulatory exposure StartNerve eliminates in under a minute."
        },
        {
            "num": "02",
            "title": "THE COMPLIANCE PROBLEM",
            "subtitle": "UNSCREENED IMPURITIES ARE A MATURING CORPORATE LIABILITY",
            "left_head": "REGULATORY EXPOSURE RULES",
            "left_body": "Unforeseen mutagenic impurities under international ICH M7 rules create multi-million dollar liability exposure and strict pipeline penalties for generic manufacturers.",
            "right_head": "THE WET-LAB BOTTLENECK",
            "right_body": "Physical Ames testing consumes massive operational budgets and weeks of pipeline time. At industrial throughput levels, wet-lab-only screening is a structural bottleneck.",
            "notes": "Traditional Ames testing is applied blanket-wide regardless of risk. Titan V11 triages the stream instantly at a marginal computational cost."
        },
        {
            "num": "03",
            "title": "MARKET TIMING INDICATORS",
            "subtitle": "THE ENFORCEMENT GATE IS TIGHTENING FASTER THAN FACTORY ADAPTATION",
            "left_head": "HIGH THROUGHPUT / THIN MARGINS",
            "left_body": "Generic manufacturers operate on aggressive schedules and thin operational margins. Global regulatory scrutiny is rising, meaning companies can no longer absorb the cost of compliance failures.",
            "right_head": "THE DEPLOYMENT GAP",
            "right_body": "Existing computational tools are built for academic research or early-stage drug discovery—completely missing the industrial generic-manufacturing compliance workflow.",
            "notes": "The regulatory net is tightening exactly as generic manufacturing volume scales—that gap is our market."
        },
        {
            "num": "04",
            "title": "THE DUAL-STREAM TECHNOLOGY",
            "subtitle": "GEOMETRY-AWARE DEEP LEARNING ARCHITECTURE OPERATING ON A SINGLE MOLECULE",
            "left_head": "STREAM 1: 2D TOPOLOGY GRAPH",
            "left_body": "GATv2 Graph Attention layers map 162-dimensional atom-level feature vectors cleanly across the electronic connection topology matrix.",
            "right_head": "STREAM 2: 3D SPATIAL GEOMETRY",
            "right_body": "SchNet physical encoders compute true Angstrom-scale spatial coordinates via local ETKDGv3 conformer generation and MMFF94 force-field energy minimization.",
            "notes": "Electronic topology and physical geometry are fused via gated linear projections into one unified 12-pathway target matrix profile simultaneously."
        },
        {
            "num": "05",
            "title": "THE CORE DATA MOAT",
            "subtitle": "58,013 VERIFIED 3D MOLECULAR GRAPHS — REGULATORY-GRADE INFRASTRUCTURE",
            "left_head": "GROUND TRUTH PIPELINE",
            "left_body": "Sourced entirely from federal-grade EPA invitrodb v4.3 and DSSTox data networks. This represents high-fidelity, regulatory-grade toxicology data assets, not scraped chemistry strings.",
            "right_head": "GEOMETRIC RESOLUTION",
            "right_body": "Every single molecule is completely geometry-resolved on disk in 3D space, heavily augmented via systematic SMILES text enumeration to expand boundary coverage.",
            "notes": "This is a clean, geometry-verified data asset that took real engineering time to compile, and it compounds with every single client screen."
        },
        {
            "num": "06",
            "title": "VALIDATION BY SCAFFOLD SPLITS",
            "subtitle": "HONEST industry BENCHMARKS vs. INFLATED RANDOM CHEMTYPE OVERLAPS",
            "left_head": "BEMIS-MURCKO RIGOR",
            "left_body": "Evaluated strictly under held-out scaffold splits where testing molecules share no core structures with training segments. It tests true generalization, not surface memorization.",
            "right_head": "THE P53 SCORE ADVANTAGE",
            "right_body": "Achieving a verified 0.8603 ROC-AUC score on the highly critical p53 tumor-suppressor pathway. Bypasses the random-split shortcuts commonly used by academic tools.",
            "notes": "We validated the hard way—scaffold splits that punish memorization—because our clients are betting real compliance decisions on this score."
        },
        {
            "num": "07",
            "title": "TRIAGING THE WORKFLOW",
            "subtitle": "UPSTREAM OPERATIONAL FILTERING BEFORE COSTLY WET-LAB ASSAYS",
            "left_head": "THE FIREWALL PIPELINE",
            "left_body": "Titan V11 sits directly upstream of physical pipelines, instantly issuing a Green, Amber, or Red risk flag per pathway in under a minute.",
            "right_head": "CAPITAL EFFICIENCY",
            "right_body": "Corporate wet-lab budgets concentrate exclusively on genuine ambiguities or high-risk flags, saving weeks of unnecessary blanket-wide physical testing schedules.",
            "notes": "We accelerate compliance velocity. We do not replace regulatory-mandated physical testing protocols."
        },
        {
            "num": "08",
            "title": "COMMERCIAL POSITIONING MATRIX",
            "subtitle": "DECISION-SUPPORT INFRASTRUCTURE BOUNDED FOR PROCUREMENT APPROVAL",
            "left_head": "REJECTED DIAGNOSTIC CLAIMS",
            "left_body": "Implies a legal certainty a predictive model cannot claim, inviting open-ended corporate liability exposure and dragging deals into 6-month legal reviews.",
            "right_head": "BOUNDED DECISION-SUPPORT",
            "right_body": "Structured as a Software decision-support tool with clear Limitation-of-Liability clauses inside every SLA, allowing enterprise procurement layers to approve fast.",
            "notes": "We built the legal architecture first—this is a tool procurement teams can approve without a six-month liability review."
        },
        {
            "num": "09",
            "title": "GO-TO-MARKET AND TRACK SCALING",
            "subtitle": "BEACHHEAD GEOGRAPHIC CORRIDORS AND TIERED SUBSCRIPTION MODULES",
            "left_head": "INDUSTRIAL LANDING HUBS",
            "left_body": "Directly targeting high-volume bulk drug manufacturing clusters and generic API factory networks across primary industrial corridors: Pune, Hyderabad, and Ahmedabad.",
            "right_head": "REVENUE STEP PIPELINE",
            "right_body": "Transitioning from individual pilot validation batch audits into recurring multi-line pipeline licenses to target a clear 10 Lakh revenue run-rate milestone by November 2026.",
            "notes": "We're not selling software into a research lab—we're selling compliance insurance into a factory's daily throughput."
        },
        {
            "num": "10",
            "title": "THE STRATEGIC ALIGNMENT ASK",
            "subtitle": "SCALING PROPRIETARY DRY-LAB FIREWALL INFRASTRUCTURE ACROSS THE MARKET",
            "left_head": "COMPETITION ADVANCEMENT",
            "left_body": "Securing track recognition and ecosystem placement to unlock deep corporate network visibility, industry access, and strategic mentorship assets.",
            "right_head": "COMMERCIAL PILOT VELOCITY",
            "right_body": "Securing 3 strategic pilot validation audit partnerships with bulk API manufacturers to cross-benchmark Titan V11 analytics directly against wet-lab compliance logs.",
            "notes": "The core mathematical stack is validated and operational. Our ask is strictly for market velocity. Help us scale the firewall."
        }
    ]

    for data in slides_content:
        slide = prs.slides.add_slide(blank_layout)
        
        # Premium Minimalist Warm Alabaster Base Paint
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR_BG

        # Left Column Base Box Container
        add_clean_container(slide, Inches(0.75), Inches(2.2), Inches(5.6), Inches(4.3), COLOR_NAVY, COLOR_NAVY)
        # Right Column Base Box Container
        add_clean_container(slide, Inches(6.983), Inches(2.2), Inches(5.6), Inches(4.3), COLOR_BG, COLOR_NAVY)

        # 1. Structural Header Metric Tracker (Top Margin Index)
        num_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(1.5), Inches(0.4))
        apply_text_styling(num_box.text_frame.paragraphs[0], f"SN // {data['num']}", 12, COLOR_COPPER, bold=True)

        # 2. Main Executive Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.75), Inches(11.8), Inches(0.6))
        apply_text_styling(title_box.text_frame.paragraphs[0], data["title"], 26, COLOR_NAVY, bold=True)

        # 3. Clean Underline Label Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.25), Inches(11.8), Inches(0.4))
        apply_text_styling(sub_box.text_frame.paragraphs[0], data["subtitle"], 10, COLOR_MUTED, bold=True)

        # 4. Left Column Content (High Contrast Inside Navy Block)
        l_title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.1), Inches(0.5))
        apply_text_styling(l_title_box.text_frame.paragraphs[0], data["left_head"], 15, COLOR_COPPER, bold=True)
        
        l_body_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(5.1), Inches(3.2))
        tf_lb = l_body_box.text_frame
        tf_lb.word_wrap = True
        apply_text_styling(tf_lb.paragraphs[0], data["left_body"], 14, COLOR_WHITE, bold=False)

        # 5. Right Column Content (Minimalist Deep Slate on Alabaster Box)
        r_title_box = slide.shapes.add_textbox(Inches(7.233), Inches(2.5), Inches(5.1), Inches(0.5))
        apply_text_styling(r_title_box.text_frame.paragraphs[0], data["right_head"], 15, COLOR_COPPER, bold=True)
        
        r_body_box = slide.shapes.add_textbox(Inches(7.233), Inches(3.1), Inches(5.1), Inches(3.2))
        tf_rb = r_body_box.text_frame
        tf_rb.word_wrap = True
        apply_text_styling(tf_rb.paragraphs[0], data["right_body"], 14, COLOR_NAVY, bold=False)

        # 6. Bottom Brand Tag Footnote
        foot_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.9), Inches(11.8), Inches(0.3))
        apply_text_styling(foot_box.text_frame.paragraphs[0], "STARTNERVE INTELLIGENCE   ·   TITAN DEEP LEARNING COMPLIANCE ENGINE", 9, COLOR_MUTED, bold=False)

        # --- Presenter Script Mapping ---
        # Seamless mapping of Claude's presentation scripts straight into your PowerPoint notes field
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = f"VERBAL ANCHOR:\n\"{data['notes']}\""

    out_file = "StartNerve_Alabaster_Prestige_Deck.pptx"
    prs.save(out_file)
    print("\n" + "="*80)
    print("🏁 PRESTIGE EDITORIAL DECK SUCCESSFULLY GENERATED")
    print(f"📊 Output Destination → {os.path.abspath(out_file)}")
    print("💡 Status: Text boxes balanced. Presenter scripts embedded into Slide Notes.")
    print("="*80 + "\n")

if __name__ == "__main__":
    create_editorial_deck()